import os
import logfire
from docx import Document
from pptx import Presentation


def parse_office(file_path: str):
    """
    Parses .docx and .pptx files without using unstructured.
    """
    with logfire.span("📄 Office Document Parsing", filename=file_path):
        try:
            ext = os.path.splitext(file_path)[1].lower()

            if ext == ".docx":
                text = _parse_docx(file_path)

            elif ext == ".pptx":
                text = _parse_pptx(file_path)

            else:
                raise ValueError(f"Unsupported file type: {ext}")

            if not text.strip():
                logfire.warning(f"⚠️ No text found in {file_path}")
            else:
                logfire.info(f"✅ Successfully parsed {len(text)} characters")

            return text

        except Exception as e:
            logfire.error(f"❌ Office Parse Failed: {e}")
            raise


def _parse_docx(file_path: str) -> str:
    doc = Document(file_path)

    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # Extract text from tables as well
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    paragraphs.append(cell.text)

    return "\n".join(paragraphs)


def _parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)

    slides_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slides_text.append(text)

    return "\n".join(slides_text)

